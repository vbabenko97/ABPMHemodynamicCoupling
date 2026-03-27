#!/usr/bin/env python3
"""Generate IEEE ELNANO 2026 presentation.

Detecting Stress-Linked Blood Pressure Relationship Shifts
from Ambulatory Monitoring.

Usage:
    python generate_presentation.py

Output:
    presentation.pptx in the current working directory.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Color / style constants
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY = RGBColor(0x0F, 0x1A, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF6, 0xF8)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
BLUE_ACCENT = RGBColor(0x2E, 0x86, 0xAB)
RED_ACCENT = RGBColor(0xE0, 0x52, 0x63)
MUTED_GRAY = RGBColor(0x88, 0x88, 0x88)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
VERY_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFA)
PIPELINE_BG = RGBColor(0xE1, 0xEB, 0xF5)

FONT_BODY = "Calibri"
FONT_HEADER = "Calibri"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

HEADER_HEIGHT = Inches(1.0)
FOOTER_HEIGHT = Inches(0.4)
CONTENT_TOP = Inches(1.25)
CONTENT_LEFT = Inches(0.7)
CONTENT_WIDTH = Inches(11.9)
CONTENT_HEIGHT = Inches(5.5)

# Figure paths (relative to script directory)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
FIG1_PRESENTATION_PATH = os.path.join(
    SCRIPT_DIR, "results", "presentation_cognitive_patterns.png",
)
FIG1_PATH = os.path.join(SCRIPT_DIR, "results", "fig1.png")
CONTEXT_BREAKDOWN_PATH = os.path.join(
    SCRIPT_DIR, "results", "presentation_context_breakdown.png",
)
LOGO_PATH = os.path.join(SCRIPT_DIR, "logo", "abpm_analysis.png")


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def create_presentation():
    """Create and configure the base presentation object."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs):
    """Add a blank slide with light gray background."""
    layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY_BG
    return slide


def add_header_bar(slide, title_text):
    """Add a dark navy header bar with white title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Thin accent line at bottom of header
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), HEADER_HEIGHT, SLIDE_WIDTH, Inches(0.04),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE_ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.15), Inches(11.9), Inches(0.7),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.LEFT


def add_footer(slide, slide_number=None):
    """Add IEEE ELNANO 2026 footer and optional slide number."""
    # Footer text
    footer_box = slide.shapes.add_textbox(
        Inches(0.7),
        SLIDE_HEIGHT - Inches(0.5),
        Inches(4),
        Inches(0.35),
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IEEE ELNANO 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED_GRAY
    p.font.name = FONT_BODY
    p.font.italic = True

    # Slide number
    if slide_number is not None:
        num_box = slide.shapes.add_textbox(
            Inches(12.1),
            SLIDE_HEIGHT - Inches(0.5),
            Inches(0.9),
            Inches(0.35),
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_number)
        p.font.size = Pt(9)
        p.font.color.rgb = MUTED_GRAY
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.RIGHT


def add_speaker_notes(slide, notes_text):
    """Set the speaker notes for a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_bullet_list(slide, bullets, left=None, top=None, width=None, height=None,
                    font_size=16, color=CHARCOAL, line_spacing=1.4, bold_first=False):
    """Add a text box with bullet points."""
    left = left or CONTENT_LEFT
    top = top or CONTENT_TOP
    width = width or CONTENT_WIDTH
    height = height or CONTENT_HEIGHT

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle bold prefix (text before colon)
        if bold_first and ":" in bullet_text:
            colon_idx = bullet_text.index(":")
            bold_part = bullet_text[: colon_idx + 1]
            rest_part = bullet_text[colon_idx + 1 :]
            run_bold = p.add_run()
            run_bold.text = bold_part
            run_bold.font.bold = True
            run_bold.font.size = Pt(font_size)
            run_bold.font.color.rgb = color
            run_bold.font.name = FONT_BODY
            run_rest = p.add_run()
            run_rest.text = rest_part
            run_rest.font.size = Pt(font_size)
            run_rest.font.color.rgb = color
            run_rest.font.name = FONT_BODY
        else:
            p.text = bullet_text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = FONT_BODY

        p.space_after = Pt(8)
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u2022")
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgb = etree.SubElement(buClr, qn("a:srgbClr"))
        srgb.set("val", "2E86AB")


def add_stat_callout(slide, number_text, label_text, x, y, w=2.5, h=1.5,
                     number_color=NAVY, bg_color=WHITE):
    """Add a large-number stat callout card."""
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(0.75)

    # Number
    num_box = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(0.8),
    )
    tf = num_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = number_color
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.85), Inches(w - 0.3), Inches(0.5),
    )
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label_text
    p.font.size = Pt(11)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER


def add_pipeline_box(slide, text, x, y, w, h, fill_color=WHITE, text_color=CHARCOAL,
                     font_size=11, bold=False):
    """Add a rounded rectangle box for the pipeline diagram."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = BLUE_ACCENT
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = FONT_BODY
    p.font.bold = bold
    # Vertical center
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def add_arrow(slide, x, y, w):
    """Add a horizontal arrow connector."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.25),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE_ACCENT
    arrow.line.fill.background()


def try_add_image(slide, path, x, y, w, h, placeholder_text="[Figure not found]"):
    """Try to insert an image; fall back to placeholder text box."""
    try:
        if os.path.isfile(path):
            slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
            return True
        else:
            raise FileNotFoundError(f"File not found: {path}")
    except Exception:
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = placeholder_text
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED_GRAY
        p.font.name = FONT_BODY
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
        return False


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_slide_01_title(prs):
    """Slide 1: Title slide -- dark background, centered content."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY

    # Thin accent stripe near top
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), SLIDE_WIDTH, Inches(0.04),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = BLUE_ACCENT
    stripe.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.4),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Detecting Stress-Linked Blood Pressure Relationship Shifts from Ambulatory Monitoring"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Authors
    authors_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.7),
    )
    tf = authors_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "A. Tymchak, A. Butkevych, K. Yurtsiv, I. Nastenko, "
        "D. Zahorska, O. Shaposhnyk, S. Yanushkevich, V. Babenko"
    )
    p.font.size = Pt(16)
    p.font.color.rgb = ICE_BLUE
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Affiliations
    aff_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.9),
    )
    tf = aff_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Igor Sikorsky Kyiv Polytechnic Institute  |  "
        "University of Calgary  |  "
        "Center of Maternity & Childhood, NAMS of Ukraine"
    )
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x99, 0xAA, 0xCC)
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Conference
    conf_box = slide.shapes.add_textbox(
        Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.5),
    )
    tf = conf_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IEEE ELNANO 2026  \u2022  Kyiv, Ukraine"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # NATO acknowledgment
    nato_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(5.7), Inches(8.3), Inches(0.5),
    )
    tf = nato_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Supported by NATO Science for Peace and Security Programme \u2013 Grant G8475"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x88, 0x99, 0xBB)
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Bottom stripe
    stripe2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.04),
    )
    stripe2.fill.solid()
    stripe2.fill.fore_color.rgb = BLUE_ACCENT
    stripe2.line.fill.background()

    # No slide number on title slide; add only the IEEE footer
    footer_box = slide.shapes.add_textbox(
        Inches(5.0), Inches(6.95), Inches(3.3), Inches(0.35),
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IEEE ELNANO 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x66, 0x77, 0x99)
    p.font.name = FONT_BODY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, (
        "Good morning, everyone. My name is Vitalii Babenko, presenting on behalf of our "
        "research team from Igor Sikorsky KPI, the University of Calgary, and the Center of "
        "Maternity and Childhood in Kyiv. Today, we will discuss a computational method "
        "designed to detect stress-linked shifts in the relationship between blood pressure "
        "and heart rate, utilizing standard ambulatory monitoring data. I would like to "
        "gratefully acknowledge the NATO Science for Peace and Security Programme for "
        "supporting this interdisciplinary research. We will examine how shifting from "
        "cohort averages to individualized baseline modeling can help us extract clearer, "
        "contextual stress signals from everyday clinical recordings."
    ))
    return slide


def build_slide_02_motivation(prs):
    """Slide 2: Motivation / Problem."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Motivation: The Specificity Problem")

    bullets = [
        "Stress impacts cardiovascular regulation and blood pressure",
        "Specificity gap in standard Ambulatory Blood Pressure Monitoring (ABPM)",
        "Elevated readings: Acute stress response vs. sustained hypertension",
        "Real-world setting: Kyiv residents facing air-raid alert exposures",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), height=Inches(3.0),
                    font_size=17, bold_first=False)

    # Visual element: key insight callout box (moved up to reduce gap)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.0),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = VERY_LIGHT_BLUE
    box.line.color.rgb = BLUE_ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key question: "
    run.font.bold = True
    run.font.size = Pt(15)
    run.font.color.rgb = NAVY
    run.font.name = FONT_BODY
    run2 = p.add_run()
    run2.text = "Can standard 24-hour ABPM data reveal the physiological footprint of different contexts through changes in hemodynamic coupling?"
    run2.font.size = Pt(15)
    run2.font.color.rgb = CHARCOAL
    run2.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 2)
    add_speaker_notes(slide, (
        "The primary clinical challenge we are addressing is a lack of specificity in standard "
        "cardiovascular monitoring. A 24-hour Ambulatory Blood Pressure Monitor, or ABPM, "
        "successfully captures dynamic changes in blood pressure, but it fundamentally lacks "
        "behavioral context. When a clinician reviews a 24-hour log, an elevated reading could "
        "reflect several different clinical realities. It might be a residual white-coat effect, "
        "emerging sustained hypertension that requires pharmacological intervention, or simply a "
        "normal, acute autonomic response to a transient environmental stressor. Without context, "
        "these highly distinct physiological states look completely identical on paper. This "
        "feasibility study investigates whether we can extract that missing contextual signal "
        "directly from the ABPM data itself. Our dataset was collected in a real-world setting "
        "in Kyiv, where participants navigated their normal daily routines, which unfortunately "
        "included periods of air-raid alert exposure over the winter. We aimed to determine if "
        "these standard, 24-hour recordings could reveal the physiological footprint of different "
        "contexts through measurable changes in an individual's hemodynamic coupling."
    ))
    return slide


def build_slide_03_related_work(prs):
    """Slide 3: Related work gaps."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Related Work Gaps")

    bullets = [
        "Group-level analyses often obscure individual heterogeneity",
        "Traditional monitoring overlooks within-subject hemodynamic coupling",
        "High-fidelity stress detection typically requires specialized equipment",
        "Need for scalable insights from deployed clinical hardware",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), width=Inches(6.0),
                    height=Inches(3.5), font_size=17)

    # Right-side comparison: two stacked cards
    # Card 1: Existing approaches
    card1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(1.5), Inches(4.8), Inches(1.8),
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE
    card1.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card1.line.width = Pt(0.75)

    # Left accent
    acc1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(1.5), Inches(0.08), Inches(1.8),
    )
    acc1.fill.solid()
    acc1.fill.fore_color.rgb = RED_ACCENT
    acc1.line.fill.background()

    tb1 = slide.shapes.add_textbox(Inches(8.1), Inches(1.6), Inches(4.3), Inches(1.6))
    tf = tb1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Existing approaches"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.font.name = FONT_BODY
    p2 = tf.add_paragraph()
    p2.text = "Cohort averages, lab-only sensors, HR/HRV confounded by exertion"
    p2.font.size = Pt(12)
    p2.font.color.rgb = CHARCOAL
    p2.font.name = FONT_BODY
    p2.space_before = Pt(6)

    # Card 2: Our approach
    card2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(3.6), Inches(4.8), Inches(1.8),
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card2.line.width = Pt(0.75)

    acc2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(3.6), Inches(0.08), Inches(1.8),
    )
    acc2.fill.solid()
    acc2.fill.fore_color.rgb = BLUE_ACCENT
    acc2.line.fill.background()

    tb2 = slide.shapes.add_textbox(Inches(8.1), Inches(3.7), Inches(4.3), Inches(1.6))
    tf = tb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Our approach"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = FONT_BODY
    p2 = tf.add_paragraph()
    p2.text = "Within-subject baselines, ABPM-only, software-level coupling deviation screening"
    p2.font.size = Pt(12)
    p2.font.color.rgb = CHARCOAL
    p2.font.name = FONT_BODY
    p2.space_before = Pt(6)

    add_footer(slide, 3)
    add_speaker_notes(slide, (
        "Much of the existing literature regarding cardiovascular stress response relies heavily "
        "on group-level statistical designs, measuring the cohort's average response to a given "
        "stimulus. However, hemodynamic responses to stress are highly heterogeneous. One patient "
        "might react to a stressor primarily with heart rate elevation, while another might "
        "experience profound vasoconstriction with a minimal chronotropic response. Averaging "
        "these reactions across a broad cohort frequently washes out these individualized "
        "physiological adjustments entirely. Furthermore, the studies that do successfully map "
        "out detailed, individual hemodynamic pathways typically utilize specialized, "
        "non-ambulatory equipment such as continuous beat-to-beat finger cuffs or lab-based "
        "impedance cardiography. While valuable, these setups cannot be scaled for prolonged "
        "at-home monitoring. Our objective is to bridge this gap by developing a scalable, "
        "software-level approach that utilizes the widely-available ABPM hardware already "
        "deployed in clinical practice. Instead of comparing patients to a population mean, we "
        "focus on interpreting the internal coupling relationships between a single patient's "
        "own cardiovascular variables over time."
    ))
    return slide


def build_slide_04_data(prs):
    """Slide 4: Data."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Data Collection & Labeling")

    bullets = [
        "28 participants (age 17\u201366; 42.9% female), Kyiv, Nov 2023 \u2013 Feb 2024",
        "2,164 valid readings across 5 labeled context windows",
        "8-level deterministic label priority; air-raid alerts take precedence",
        "Cognitive tasks: Stroop, reaction time, visuomotor coordination",
        "Physical task: modified Harvard Step Test",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), width=Inches(5.6),
                    height=Inches(3.5), font_size=17)

    # Right side: context breakdown bar chart — wider, shifted left for more bar room
    try_add_image(
        slide, CONTEXT_BREAKDOWN_PATH,
        x=6.1, y=1.3, w=6.7, h=3.8,
        placeholder_text="[Dataset composition by labeled context]",
    )

    # Bottom note
    note_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(5.6), Inches(6.5), Inches(0.4),
    )
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Full condition-level summary in Table 1 of the paper"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED_GRAY
    p.font.name = FONT_BODY
    p.font.italic = True

    add_footer(slide, 4)
    add_speaker_notes(slide, (
        "Our dataset was collected continuously between November 2023 and February 2024. We "
        "outfitted 28 participants in Kyiv with standard ABPM devices. Following artifact "
        "removal, the study yielded 2,164 valid readings. To create distinct context windows "
        "for our modeling, participants were asked to perform specific controlled activities. "
        "For cognitive stress, this included tasks like the Stroop test and visuomotor tracking; "
        "for physical stress, participants performed the Harvard Step Test. We also integrated "
        "timestamps from the Ukraine Alarm API to flag air-raid alert exposure windows. To "
        "manage overlapping contexts, we implemented an 8-level deterministic priority "
        "hierarchy. Safety alerts were prioritized at the top -- meaning any reading during an "
        "air-raid alert was exclusively labeled as such -- while baseline resting states formed "
        "the bottom of the hierarchy."
    ))
    return slide


def build_slide_05_methods(prs):
    """Slide 5: Methods -- includes pipeline diagram."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Methods: Individualized Coupling Model")

    # Full-width bullet list with formulas folded in
    bullets = [
        "Per-subject baseline: predicted DBP = f(SBP, HR)",
        "Compact feature library with interaction and inverse terms",
        "3-fold contiguous time-block cross-validation (no temporal leakage)",
        "MAE Inflation: how much worse is prediction under stress?",
        "Signed Residual Bias: is observed DBP above or below prediction?",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), width=Inches(11.9),
                    height=Inches(2.6), font_size=18)

    # Pipeline diagram — taller and more prominent
    pipeline_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.6),
    )
    pipeline_bg.fill.solid()
    pipeline_bg.fill.fore_color.rgb = PIPELINE_BG
    pipeline_bg.line.fill.background()

    # Pipeline label
    plabel = slide.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(3.0), Inches(0.35))
    tf = plabel.text_frame
    p = tf.paragraphs[0]
    p.text = "Analysis Pipeline"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY

    # Pipeline boxes (larger for projection readability)
    box_y = 4.9
    box_h = 1.0
    box_w = 1.7
    gap = 0.15
    arrow_w = 0.3

    steps = [
        "ABPM\nData",
        "Context\nLabeling",
        "Baseline\nModel Fit",
        "Time-Block\nCross-Val",
        "Deviation\nMetrics",
        "Screening\nFlags",
    ]

    start_x = 0.7
    for i, step_text in enumerate(steps):
        bx = start_x + i * (box_w + gap + arrow_w)
        add_pipeline_box(slide, step_text, bx, box_y, box_w, box_h,
                         fill_color=WHITE, text_color=NAVY, font_size=13, bold=True)
        if i < len(steps) - 1:
            ax = bx + box_w + 0.03
            add_arrow(slide, ax, box_y + box_h / 2 - 0.125, arrow_w)

    add_footer(slide, 5)
    add_speaker_notes(slide, (
        "Our analytical framework relies entirely on individualized baseline modeling. Rather "
        "than building a generalized cohort model, we trained a distinct regression model for "
        "each participant using only their non-stressful awake baseline readings. The model "
        "predicts Diastolic Blood Pressure using concurrent values of Systolic Blood Pressure "
        "and Heart Rate. To ensure the models can capture non-linear cardiovascular "
        "relationships without overfitting on sparse data, we defined a compact 6-feature "
        "library comprising the raw variables, their mathematical inverses, and multi-variable "
        "interaction terms. The candidate models were evaluated using three-fold contiguous "
        "time-block cross-validation. This contiguous splitting strategy is critical, as it "
        "prevents temporal leakage between adjacent blood pressure readings. Ultimately, simple "
        "Ordinary Least Squares regression yielded the lowest error for the majority of "
        "participants and was selected as the baseline architecture. During the labeled context "
        "windows, we quantified deviations from this baseline using two primary metrics. First, "
        "Mean Absolute Error Inflation measures the magnitude of the model's prediction "
        "degradation relative to baseline. Second, Signed Residual Bias measures the direction "
        "of the error, indicating whether the true measured diastolic pressure is systematically "
        "higher or systematically lower than what the baseline coupling relationship predicts."
    ))
    return slide


def build_slide_06_results(prs):
    """Slide 6: Results -- with Figure 1."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Results: Individual-Level Patterns")

    # Prominent null-result banner across the top
    null_banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.6),
    )
    null_banner.fill.solid()
    null_banner.fill.fore_color.rgb = RGBColor(0xFD, 0xEE, 0xEF)  # light red tint
    null_banner.line.color.rgb = RED_ACCENT
    null_banner.line.width = Pt(1.5)
    tf = null_banner.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Cohort-level contrasts not significant after BH-FDR correction (all q > 0.05)"
    run1.font.size = Pt(15)
    run1.font.bold = True
    run1.font.color.rgb = RED_ACCENT
    run1.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Key stats — shifted down to sit below the banner
    add_stat_callout(slide, "6.89", "Baseline MAE\n(mmHg)", 0.7, 2.05, w=2.8, h=1.1,
                     number_color=NAVY)
    add_stat_callout(slide, "+4.18", "Cognitive Bias\n(mmHg)", 3.8, 2.05, w=2.8, h=1.1,
                     number_color=BLUE_ACCENT)
    add_stat_callout(slide, "-2.01", "Physical Bias\n(mmHg)", 6.9, 2.05, w=2.8, h=1.1,
                     number_color=RED_ACCENT)
    add_stat_callout(slide, "20/27", "Flagged (74.1%)\n(descriptive threshold)", 10.0, 2.05, w=2.8, h=1.1,
                     number_color=NAVY)

    # Figure 1 -- insert the dot plot
    # Expected: results/fig1.png (two-panel dot plot)
    figure_path = FIG1_PRESENTATION_PATH if os.path.exists(FIG1_PRESENTATION_PATH) else FIG1_PATH
    fig_found = try_add_image(
        slide, figure_path,
        x=1.5, y=3.3, w=10.3, h=3.5,
        placeholder_text=(
            "[Figure 1: Two-panel dot plot showing MAE inflation (Panel A) "
            "and signed residual bias (Panel B) per participant]"
        ),
    )

    # Caption
    cap_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.35),
    )
    tf = cap_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Figure 1: Per-participant cognitive-task MAE inflation (A) "
        "and signed residual bias (B). "
        "Open circles = not flagged; crosses = flagged."
    )
    p.font.size = Pt(10)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_BODY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 6)
    add_speaker_notes(slide, (
        "Under baseline resting conditions, the individualized models successfully predicted "
        "diastolic pressure with a median absolute error of 6.89 millimeters of mercury. When "
        "we conducted standard cohort-level statistical tests across the context labels, the "
        "results were entirely null following Benjamini-Hochberg false discovery rate correction. "
        "At the population level, the physiological effects were statistically insignificant. "
        "However, examining the residuals at the participant level reveals distinct underlying "
        "patterns. If you refer to Figure 1 on the screen, this two-panel dot plot illustrates "
        "the MAE inflation and signed bias for each subject. We applied a descriptive flagging "
        "threshold, identifying participants whose prediction error inflated by over 50 percent, "
        "or whose bias shifted by more than 2 millimeters of mercury. Based on these criteria, "
        "20 of 27 participants met the descriptive flagging rule. Notably, the directional bias "
        "varied systematically according to the stressor type. During cognitive tasks, the models "
        "consistently under-predicted diastolic pressure, yielding a positive median bias of "
        "4.18 millimeters of mercury. Conversely, during physical exertion, we observed a "
        "negative median bias of 2.01 millimeters of mercury. Finally, in the sub-cohort of "
        "11 participants exposed to air-raid alerts, we recorded substantial prediction error "
        "inflation averaging 27.35 percent, alongside a near-zero directional bias."
    ))
    return slide


def build_slide_07_discussion(prs):
    """Slide 7: Discussion."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Discussion: Interpreting the Patterns")

    # Two-column layout
    # Left column: bullets
    bullets = [
        "Null group finding supports within-subject analytical approach",
        "Divergent biases suggest context-dependent coupling shifts",
        "Residual analysis serves as a macroscopic screening signal",
        "Air-raid condition frames heterogeneous, high-variance exposure windows",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), width=Inches(6.0),
                    height=Inches(3.0), font_size=16)

    # Right column: interpretation cards
    cards = [
        ("Cognitive tasks", "+4.18 mmHg bias", "Consistent with increased\nvascular tone (inferential)", BLUE_ACCENT),
        ("Physical tasks", "-2.01 mmHg bias", "Consistent with exercise-related\nperipheral vasodilation (inferential)", RED_ACCENT),
        ("Air-raid alerts", "+27.35% MAE inflation", "Heterogeneous behavioral\nresponses, near-zero bias", NAVY),
    ]

    for idx, (title, stat, desc, accent_color) in enumerate(cards):
        cy = 1.5 + idx * 1.55  # Even spacing between cards
        # Card bg
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7.5), Inches(cy), Inches(5.1), Inches(1.4),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        card.line.width = Pt(0.75)

        # Left accent
        acc = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7.5), Inches(cy), Inches(0.08), Inches(1.4),
        )
        acc.fill.solid()
        acc.fill.fore_color.rgb = accent_color
        acc.line.fill.background()

        # Title + stat
        tb = slide.shapes.add_textbox(Inches(7.85), Inches(cy + 0.1), Inches(4.5), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = title + "  "
        run1.font.size = Pt(13)
        run1.font.bold = True
        run1.font.color.rgb = accent_color
        run1.font.name = FONT_BODY
        run2 = p.add_run()
        run2.text = stat
        run2.font.size = Pt(13)
        run2.font.color.rgb = CHARCOAL
        run2.font.name = FONT_BODY

        # Description
        db = slide.shapes.add_textbox(Inches(7.85), Inches(cy + 0.55), Inches(4.5), Inches(0.7))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = FONT_BODY

    add_footer(slide, 7)
    add_speaker_notes(slide, (
        "The contrast between our null cohort-level findings and the observable individual-level "
        "deviations strongly supports our core premise: relying on group-averaging may obscure "
        "meaningful within-subject physiological adjustments. The directional split observed "
        "between tasks aligns with expected patterns of peripheral vascular resistance. Focused "
        "cognitive tasks are generally associated with sympathetic vasoconstriction, which may "
        "drive diastolic pressure higher relative to concurrent heart rate. In contrast, heavy "
        "physical exertion typically prompts muscular vasodilation, lowering peripheral "
        "resistance and creating a negative diastolic bias relative to the baseline projection. "
        "Therefore, our residual metrics may serve as an indirect screening signal for these "
        "types of vascular adjustments. Regarding the air-raid condition, characterizing these "
        "as exposure windows rather than direct measures of acute stress is an important "
        "distinction. The combination of high error inflation and near-zero directional bias "
        "suggests a highly heterogeneous response. This variance likely reflects differing "
        "behavioral reactions to the alerts -- ranging from physical relocation to a shelter, "
        "to remaining sedentary at a desk."
    ))
    return slide


def build_slide_08_limitations(prs):
    """Slide 8: Limitations and Future Work."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Limitations & Future Directions")

    # Two-column layout
    # Left: Limitations
    lim_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.4),
    )
    tf = lim_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Current Limitations"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.font.name = FONT_BODY

    lim_bullets = [
        "Sample size (N=28) and sparse ABPM cadence (15-30 min)",
        "Lack of continuous cardiac output ground truth",
        "Coarse context labels do not capture moment-to-moment arousal",
    ]
    add_bullet_list(slide, lim_bullets, left=Inches(0.7), top=Inches(1.85),
                    width=Inches(5.5), height=Inches(2.5), font_size=15)

    # Right: Future Work
    fut_title = slide.shapes.add_textbox(
        Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.4),
    )
    tf = fut_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Future Directions"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = FONT_BODY

    fut_bullets = [
        "Larger cohorts with higher-resolution ambulatory monitors",
        "Impedance cardiography calibration sub-study",
        "Integration with cuffless / wearable BP platforms",
        "Ecological momentary assessment prompts at cuff inflation",
    ]
    add_bullet_list(slide, fut_bullets, left=Inches(7.0), top=Inches(1.85),
                    width=Inches(5.8), height=Inches(3.0), font_size=15)

    # Visual separator line (more visible)
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.55), Inches(1.5), Inches(0.04), Inches(3.5),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
    sep.line.fill.background()

    # Bottom callout (more generous margins from slide edges)
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.8), Inches(5.2), Inches(9.7), Inches(1.0),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = VERY_LIGHT_BLUE
    callout.line.color.rgb = BLUE_ACCENT
    callout.line.width = Pt(1.5)
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "As next-generation wearable and cuffless BP technologies mature, "
        "individualized software-level screening becomes increasingly viable "
        "for continuous, context-aware cardiovascular monitoring."
    )
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 8)
    add_speaker_notes(slide, (
        "Because this is an initial feasibility study, it carries notable limitations. Our "
        "sample size of 28 participants is relatively small, which leaves subset analyses -- "
        "particularly the air-raid exposure group -- statistically underpowered. Additionally, "
        "standard ABPM's sampling cadence of 15 to 30 minutes provides a very sparse "
        "representation of human autonomic regulation, which is a continuous biological process. "
        "Finally, without simultaneous ground-truth measurements of cardiac output and systemic "
        "vascular resistance, our mechanistic interpretations regarding vasoconstriction and "
        "vasodilation remain inferential. Moving forward, our next steps include validating "
        "these analytical techniques in larger cohorts using higher-resolution ambulatory "
        "monitors paired with digital ecological momentary assessments. We are also designing "
        "a dedicated laboratory sub-study utilizing impedance cardiography, which is necessary "
        "to mathematically calibrate our residual metrics against established hemodynamic gold "
        "standards. As next-generation wearable and cuffless blood pressure technologies mature, "
        "we anticipate this individualized software approach will become an increasingly viable "
        "method for continuous, context-aware physiological screening."
    ))
    return slide


def build_slide_09_conclusions(prs):
    """Slide 9: Conclusions."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Conclusions")

    conclusions = [
        ("Standard ABPM data yields within-subject stress screening signals",
         "24-hour recordings contain context-sensitive information when analyzed at the individual level"),
        ("Individualized baselines capture context missed by cohort averages",
         "Moving from group-average to per-subject models reveals regulatory decoupling episodes"),
        ("Dual-metric approach distinguishes magnitude and direction of coupling deviation",
         "MAE inflation quantifies shift magnitude; signed bias classifies hemodynamic patterns"),
        ("Scalable, software-level analysis for behavioral cardiovascular monitoring",
         "Computationally lightweight method to increase specificity in longitudinal screening"),
    ]

    y_start = 1.5
    for idx, (headline, detail) in enumerate(conclusions):
        cy = y_start + idx * 1.3

        # Number circle
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.9), Inches(cy + 0.05), Inches(0.45), Inches(0.45),
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE_ACCENT
        circ.line.fill.background()
        tf = circ.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(idx + 1)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_HEADER
        p.alignment = PP_ALIGN.CENTER

        # Headline
        hb = slide.shapes.add_textbox(
            Inches(1.6), Inches(cy), Inches(10.5), Inches(0.4),
        )
        tf = hb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = headline
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = FONT_BODY

        # Detail
        db = slide.shapes.add_textbox(
            Inches(1.6), Inches(cy + 0.42), Inches(10.5), Inches(0.5),
        )
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(13)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = FONT_BODY

    add_footer(slide, 9)
    add_speaker_notes(slide, (
        "To conclude, standard ambulatory blood pressure data contains valuable, "
        "context-sensitive characteristics when analyzed rigorously at the individual level. "
        "Moving away from generalized group-average models toward individualized baseline "
        "approaches enables us to capture discrete episodes of regulatory decoupling that "
        "would otherwise remain hidden. By employing a dual-metric assessment -- quantifying "
        "both the magnitude of the model's prediction failure and the direction of the "
        "prediction error -- we can classify deviations that suggest differing hemodynamic "
        "patterns. Ultimately, this represents a scalable, computationally lightweight method "
        "intended to increase the specificity and contextual intelligence of longitudinal "
        "cardiovascular screening."
    ))
    return slide


def build_slide_10_thank_you(prs):
    """Slide 10: Acknowledgments and Thank You."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY

    # Top accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.0), SLIDE_WIDTH, Inches(0.04),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = BLUE_ACCENT
    stripe.line.fill.background()

    # "Thank You" title
    ty_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.4), Inches(11.3), Inches(1.0),
    )
    tf = ty_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Acknowledgments
    ack_items = [
        "Funding: NATO Science for Peace and Security Programme (Grant G8475)",
        "Collaborators: Igor Sikorsky KPI & University of Calgary",
        "Special thanks to the study participants in Kyiv",
    ]
    ack_box = slide.shapes.add_textbox(
        Inches(2.0), Inches(2.8), Inches(9.3), Inches(2.2),
    )
    tf = ack_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(ack_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = ICE_BLUE
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

    # Questions box
    q_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.8),
    )
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = BLUE_ACCENT
    q_box.line.fill.background()
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Bottom stripe
    stripe2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.6), SLIDE_WIDTH, Inches(0.04),
    )
    stripe2.fill.solid()
    stripe2.fill.fore_color.rgb = BLUE_ACCENT
    stripe2.line.fill.background()

    # Footer
    footer_box = slide.shapes.add_textbox(
        Inches(5.0), Inches(6.85), Inches(3.3), Inches(0.35),
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IEEE ELNANO 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x66, 0x77, 0x99)
    p.font.name = FONT_BODY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    # Slide number
    num_box = slide.shapes.add_textbox(
        Inches(12.1), Inches(6.85), Inches(0.9), Inches(0.35),
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = "10"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x66, 0x77, 0x99)
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.RIGHT

    add_speaker_notes(slide, (
        "I would like to conclude by formally thanking the NATO Science for Peace and Security "
        "Programme under grant G8475 for funding this research. I also extend my gratitude to "
        "my collaborators, and most importantly, to our study participants for their remarkable "
        "adherence to the clinical protocol while navigating challenging environmental conditions "
        "in Kyiv. Thank you very much for your time and attention today. I would now be happy to "
        "take any questions you may have."
    ))
    return slide


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    prs = create_presentation()

    build_slide_01_title(prs)
    build_slide_02_motivation(prs)
    build_slide_03_related_work(prs)
    build_slide_04_data(prs)
    build_slide_05_methods(prs)
    build_slide_06_results(prs)
    build_slide_07_discussion(prs)
    build_slide_08_limitations(prs)
    build_slide_09_conclusions(prs)
    build_slide_10_thank_you(prs)

    output_path = os.path.join(SCRIPT_DIR, "presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
