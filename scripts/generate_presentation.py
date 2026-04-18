#!/usr/bin/env python3
"""Generate IEEE ELNANO 2026 presentation.

Detecting Stress-Linked Blood Pressure Relationship Shifts
from Ambulatory Monitoring.

Usage:
    python scripts/generate_presentation.py

Output:
    presentation.pptx in the current working directory.
"""

import os

import matplotlib.dates as mdates
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from abpm_hemodynamic_coupling.config import Columns, Config
from abpm_hemodynamic_coupling.feature_engineering import DBPFeatureExtractor
from abpm_hemodynamic_coupling.modeling import ModelTrainer

# ---------------------------------------------------------------------------
# Color / style constants
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY = RGBColor(0x0F, 0x1A, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF6, 0xF8)
CHARCOAL = RGBColor(0x33, 0x33, 0x33)
BLUE_ACCENT = RGBColor(0x2E, 0x86, 0xAB)
RED_ACCENT = RGBColor(0xE0, 0x52, 0x63)
MUTED_GRAY = RGBColor(0x88, 0x88, 0x88)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
VERY_LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFA)
PIPELINE_BG = RGBColor(0xE1, 0xEB, 0xF5)

FONT_BODY = "Calibri"
FONT_HEADER = "Calibri"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

HEADER_HEIGHT = Inches(1.0)
FOOTER_HEIGHT = Inches(0.4)
CONTENT_TOP = Inches(1.25)
CONTENT_LEFT = Inches(0.7)
CONTENT_WIDTH = Inches(11.9)
CONTENT_HEIGHT = Inches(5.5)

# Paths are resolved relative to the repository root (one level above scripts/)
SCRIPT_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG1_PRESENTATION_PATH = os.path.join(
    SCRIPT_DIR, "results", "presentation_cognitive_patterns.png",
)
FIG1_PATH = os.path.join(SCRIPT_DIR, "results", "fig1.png")
CONTEXT_BREAKDOWN_PATH = os.path.join(
    SCRIPT_DIR, "results", "presentation_context_breakdown.png",
)
LOGO_PATH = os.path.join(SCRIPT_DIR, "logo", "abpm_analysis.png")
SUBJECT_METRICS_PATHS = [
    os.path.join(SCRIPT_DIR, "docs", "thesis", "per_subject_metrics.csv"),
    os.path.join(SCRIPT_DIR, "results", "per_subject_metrics.csv"),
]
MONITORING_PATHS = [
    os.path.join(SCRIPT_DIR, "results", "thesis", "monitoring_labeled.parquet"),
    os.path.join(SCRIPT_DIR, "data", "monitoring_data.csv"),
]
CASE_STUDY_PATH = os.path.join(SCRIPT_DIR, "results", "presentation_case_study.png")
CASE_STUDY_SUBJECT = 35


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def create_presentation():
    """Create and configure the base presentation object."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs):
    """Add a blank slide with light gray background."""
    layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY_BG
    return slide


def add_header_bar(slide, title_text):
    """Add a dark navy header bar with white title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, HEADER_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Thin accent line at bottom of header
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), HEADER_HEIGHT, SLIDE_WIDTH, Inches(0.04),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE_ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.15), Inches(11.9), Inches(0.7),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.LEFT


def add_footer(slide, slide_number=None):
    """Add IEEE ELNANO 2026 footer and optional slide number."""
    # Footer text
    footer_box = slide.shapes.add_textbox(
        Inches(0.7),
        SLIDE_HEIGHT - Inches(0.5),
        Inches(4),
        Inches(0.35),
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IEEE ELNANO 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED_GRAY
    p.font.name = FONT_BODY
    p.font.italic = True

    # Slide number
    if slide_number is not None:
        num_box = slide.shapes.add_textbox(
            Inches(12.1),
            SLIDE_HEIGHT - Inches(0.5),
            Inches(0.9),
            Inches(0.35),
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(slide_number)
        p.font.size = Pt(9)
        p.font.color.rgb = MUTED_GRAY
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.RIGHT


def add_speaker_notes(slide, notes_text):
    """Set the speaker notes for a slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def add_bullet_list(slide, bullets, left=None, top=None, width=None, height=None,
                    font_size=16, color=CHARCOAL, line_spacing=1.4, bold_first=False):
    """Add a text box with bullet points."""
    left = left or CONTENT_LEFT
    top = top or CONTENT_TOP
    width = width or CONTENT_WIDTH
    height = height or CONTENT_HEIGHT

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle bold prefix (text before colon)
        if bold_first and ":" in bullet_text:
            colon_idx = bullet_text.index(":")
            bold_part = bullet_text[: colon_idx + 1]
            rest_part = bullet_text[colon_idx + 1 :]
            run_bold = p.add_run()
            run_bold.text = bold_part
            run_bold.font.bold = True
            run_bold.font.size = Pt(font_size)
            run_bold.font.color.rgb = color
            run_bold.font.name = FONT_BODY
            run_rest = p.add_run()
            run_rest.text = rest_part
            run_rest.font.size = Pt(font_size)
            run_rest.font.color.rgb = color
            run_rest.font.name = FONT_BODY
        else:
            p.text = bullet_text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = FONT_BODY

        p.space_after = Pt(8)
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u2022")
        buClr = etree.SubElement(pPr, qn("a:buClr"))
        srgb = etree.SubElement(buClr, qn("a:srgbClr"))
        srgb.set("val", "2E86AB")


def add_stat_callout(slide, number_text, label_text, x, y, w=2.5, h=1.5,
                     number_color=NAVY, bg_color=WHITE):
    """Add a large-number stat callout card."""
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(0.75)

    # Number
    num_box = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(0.8),
    )
    tf = num_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = number_text
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = number_color
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide.shapes.add_textbox(
        Inches(x + 0.15), Inches(y + 0.8), Inches(w - 0.3), Inches(0.58),
    )
    tf = lbl_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label_text
    p.font.size = Pt(12)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER


def add_formula_card(slide, title_text, formula_text, x, y, w, h, accent_color):
    """Add a compact formula card with a colored accent."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xD6, 0xDE, 0xE8)
    card.line.width = Pt(1.0)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(x + 0.18), Inches(y + 0.1), Inches(w - 0.3), Inches(0.25),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = FONT_BODY

    formula_box = slide.shapes.add_textbox(
        Inches(x + 0.18), Inches(y + 0.34), Inches(w - 0.32), Inches(h - 0.42),
    )
    tf = formula_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = formula_text
    p.font.size = Pt(16)
    p.font.color.rgb = CHARCOAL
    p.font.name = FONT_BODY


def add_pipeline_box(slide, text, x, y, w, h, fill_color=WHITE, text_color=CHARCOAL,
                     font_size=11, bold=False):
    """Add a rounded rectangle box for the pipeline diagram."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = BLUE_ACCENT
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.name = FONT_BODY
    p.font.bold = bold
    # Vertical center
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


def add_arrow(slide, x, y, w):
    """Add a horizontal arrow connector."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(0.25),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE_ACCENT
    arrow.line.fill.background()


def try_add_image(slide, path, x, y, w, h, placeholder_text="[Figure not found]"):
    """Try to insert an image; fall back to placeholder text box."""
    try:
        if os.path.isfile(path):
            slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
            return True
        raise FileNotFoundError(f"File not found: {path}")
    except Exception:
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = placeholder_text
        p.font.size = Pt(16)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER
        return False


def resolve_first_existing_path(paths):
    """Return the first existing path from a candidate list."""
    for path in paths:
        if os.path.exists(path):
            return path
    return None


def generate_presentation_cognitive_patterns():
    """Create a projector-friendly version of the cognitive dot plots."""
    metrics_path = resolve_first_existing_path(SUBJECT_METRICS_PATHS)
    if metrics_path is None:
        return False

    metrics_df = pd.read_csv(metrics_path)
    valid_df = metrics_df[metrics_df["DBP_Cognitive Task_N"] > 0].copy()
    if valid_df.empty:
        return False

    valid_df["flagged"] = (
        (valid_df["DBP_Cognitive Task_Anomaly"] > 50.0)
        | (valid_df["DBP_Cognitive Task_DeltaBias"] > 2.0)
    )

    flagged_df = valid_df[valid_df["flagged"]]
    clean_df = valid_df[~valid_df["flagged"]]
    rng = np.random.default_rng(42)

    fig, axes = plt.subplots(1, 2, figsize=(11.4, 4.8))
    panels = [
        (
            axes[0],
            "DBP_Cognitive Task_Anomaly",
            "A. MAE inflation",
            "MAE inflation (%)",
            50.0,
            "Flag threshold = 50%",
        ),
        (
            axes[1],
            "DBP_Cognitive Task_DeltaBias",
            "B. Signed residual bias",
            "Signed residual bias (mmHg)",
            2.0,
            "Flag threshold = +2 mmHg",
        ),
    ]

    for ax, column, title, ylabel, threshold, threshold_label in panels:
        x_clean = rng.normal(0.0, 0.04, size=len(clean_df))
        x_flagged = rng.normal(1.0, 0.04, size=len(flagged_df))

        if len(clean_df) > 0:
            ax.scatter(
                x_clean,
                clean_df[column],
                s=95,
                linewidth=1.5,
                marker="o",
                facecolors="none",
                edgecolors="#2E86AB",
                zorder=3,
            )
        if len(flagged_df) > 0:
            ax.scatter(
                x_flagged,
                flagged_df[column],
                s=120,
                linewidth=1.4,
                marker="X",
                color="#E05263",
                zorder=4,
            )

        median_value = float(valid_df[column].median())
        ax.axhline(threshold, color="#E05263", linewidth=3.0, alpha=0.95, zorder=2)
        ax.axhline(
            median_value,
            color="#666666",
            linestyle="--",
            linewidth=2.0,
            alpha=0.85,
            zorder=1,
        )
        ax.set_title(title, fontsize=15, fontweight="bold", pad=10)
        ax.set_ylabel(ylabel, fontsize=15, fontweight="bold")
        ax.set_xlim(-0.4, 1.4)
        ax.set_xticks([0, 1], ["Not flagged\n(n=7)", "Flagged\n(n=20)"])
        ax.tick_params(axis="x", labelsize=12)
        ax.tick_params(axis="y", labelsize=13)
        ax.grid(True, axis="y", alpha=0.25, linewidth=0.8)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)
        ax.text(
            0.02,
            0.96,
            f"Cohort median = {median_value:.2f}",
            transform=ax.transAxes,
            ha="left",
            va="top",
            fontsize=11,
            color="#666666",
            bbox=dict(boxstyle="round,pad=0.2", fc="white", ec="#D6DBDF", alpha=0.9),
        )

    fig.subplots_adjust(left=0.08, right=0.98, top=0.9, bottom=0.17, wspace=0.16)
    fig.savefig(FIG1_PRESENTATION_PATH, dpi=250, bbox_inches="tight")
    plt.close(fig)
    return True


def generate_context_breakdown_chart():
    """Create a cleaner context-composition chart for the data slide."""
    monitoring_path = resolve_first_existing_path(MONITORING_PATHS)
    if monitoring_path is None:
        return False

    if monitoring_path.endswith(".parquet"):
        monitoring_df = pd.read_parquet(monitoring_path)
    else:
        monitoring_df = pd.read_csv(monitoring_path, parse_dates=[Columns.TIME])

    categories = [
        ("Baseline", [Columns.LABEL_BASELINE], "#24355A"),
        ("Sleep", [Columns.LABEL_SLEEP], "#3B8DB3"),
        (
            "Task",
            [Columns.LABEL_COGNITIVE_TASK, Columns.LABEL_PHYSICAL_TASK],
            "#E05263",
        ),
        ("Air Alert", [Columns.LABEL_AIR_ALERT], "#6C4E9B"),
    ]

    used_mask = monitoring_df[Columns.LABEL].isin(
        [label for _, labels, _ in categories for label in labels]
    )
    data = []
    total_n = len(monitoring_df)
    for name, labels, color in categories:
        subset = monitoring_df[monitoring_df[Columns.LABEL].isin(labels)]
        data.append(
            {
                "name": name,
                "count": len(subset),
                "subjects": subset[Columns.PAT_ID].nunique(),
                "pct": len(subset) / total_n * 100.0,
                "color": color,
            }
        )

    excluded_df = monitoring_df[~used_mask]
    data.append(
        {
            "name": "Other excluded",
            "count": len(excluded_df),
            "subjects": excluded_df[Columns.PAT_ID].nunique(),
            "pct": len(excluded_df) / total_n * 100.0,
            "color": "#AAB3C2",
        }
    )

    plot_df = pd.DataFrame(data)
    fig, ax = plt.subplots(figsize=(7.7, 5.0))
    y_positions = np.arange(len(plot_df))
    bars = ax.barh(
        y_positions,
        plot_df["count"],
        color=plot_df["color"],
        height=0.56,
        edgecolor="none",
    )

    ax.set_yticks(y_positions, plot_df["name"], fontsize=13)
    ax.invert_yaxis()
    max_count = int(plot_df["count"].max())
    ax.set_xlim(0, max_count * 1.48)
    ax.set_xlabel("Valid ABPM readings", fontsize=14, fontweight="bold")
    ax.tick_params(axis="x", labelsize=11)
    ax.grid(True, axis="x", alpha=0.18, linewidth=0.8)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)

    for bar, (_, row) in zip(bars, plot_df.iterrows(), strict=False):
        y_mid = bar.get_y() + bar.get_height() / 2
        x_right = bar.get_width()
        ax.text(
            x_right + max_count * 0.03,
            y_mid - 0.02,
            f"{int(row['count']):,} readings  ({row['pct']:.1f}%)",
            va="center",
            ha="left",
            fontsize=11.6,
            fontweight="semibold",
            color="#333333",
        )
        ax.text(
            x_right + max_count * 0.03,
            y_mid + 0.23,
            f"n subjects = {int(row['subjects'])}",
            va="center",
            ha="left",
            fontsize=9.5,
            color="#666666",
        )

    fig.subplots_adjust(left=0.16, right=0.98, top=0.95, bottom=0.16)
    fig.savefig(CONTEXT_BREAKDOWN_PATH, dpi=250, bbox_inches="tight")
    plt.close(fig)
    return True


def _prepare_case_study_predictions(subject_df, winner):
    """Fit the subject-level baseline model and return predictions/residuals."""
    config = Config()
    feature_extractor = DBPFeatureExtractor(config)
    trainer = ModelTrainer(config)

    baseline_df = subject_df[subject_df[Columns.LABEL] == Columns.LABEL_BASELINE].copy()
    if baseline_df.empty:
        raise ValueError("Case-study subject has no baseline rows.")

    X_base = feature_extractor.extract(baseline_df)
    y_base = baseline_df[Columns.DBP].to_numpy(dtype=float)
    model, feature_idx, scaler = trainer.train(X_base, y_base, winner)

    X_all = feature_extractor.extract(subject_df)
    X_all_scaled = scaler.transform(X_all)
    predictions = model.predict(X_all_scaled[:, feature_idx])

    enriched = subject_df.copy()
    enriched["predicted_dbp"] = predictions
    enriched["residual_dbp"] = enriched[Columns.DBP].to_numpy(dtype=float) - predictions
    return enriched


def generate_case_study_timeline(subject_id=CASE_STUDY_SUBJECT):
    """Generate a single-subject alert timeline for the presentation."""
    metrics_path = resolve_first_existing_path(SUBJECT_METRICS_PATHS)
    monitoring_path = resolve_first_existing_path(MONITORING_PATHS)
    if metrics_path is None or monitoring_path is None:
        return False

    metrics_df = pd.read_csv(metrics_path)
    if monitoring_path.endswith(".parquet"):
        monitoring_df = pd.read_parquet(monitoring_path)
    else:
        monitoring_df = pd.read_csv(monitoring_path, parse_dates=[Columns.TIME])

    subject_metrics = metrics_df[metrics_df["participant_id"] == subject_id]
    subject_df = monitoring_df[
        monitoring_df[Columns.PAT_ID] == subject_id
    ].sort_values(Columns.TIME).copy()
    if subject_metrics.empty or subject_df.empty:
        return False

    winner = subject_metrics["DBP_Winner"].iloc[0]
    if winner == "NA":
        return False

    subject_df[Columns.TIME] = pd.to_datetime(subject_df[Columns.TIME])
    subject_df = _prepare_case_study_predictions(subject_df, winner)

    alert_df = subject_df[subject_df[Columns.LABEL] == Columns.LABEL_AIR_ALERT].copy()
    baseline_df = subject_df[subject_df[Columns.LABEL] == Columns.LABEL_BASELINE].copy()
    if alert_df.empty or baseline_df.empty:
        return False

    pre_alert_baseline = baseline_df[baseline_df[Columns.TIME] < alert_df[Columns.TIME].min()]
    if pre_alert_baseline.empty:
        pre_alert_baseline = baseline_df.iloc[: min(4, len(baseline_df))]

    display_start = pre_alert_baseline[Columns.TIME].min() - pd.Timedelta(minutes=15)
    display_end = alert_df[Columns.TIME].max() + pd.Timedelta(minutes=20)
    display_df = subject_df[
        (subject_df[Columns.TIME] >= display_start)
        & (subject_df[Columns.TIME] <= display_end)
    ].copy()
    if display_df.empty:
        return False

    alert_peak_idx = alert_df["residual_dbp"].abs().idxmax()
    alert_peak = alert_df.loc[alert_peak_idx]
    metrics_row = subject_metrics.iloc[0]

    fig = plt.figure(figsize=(11.2, 6.4))
    gs = fig.add_gridspec(3, 1, height_ratios=[1.0, 1.25, 0.85], hspace=0.12)
    ax_top = fig.add_subplot(gs[0])
    ax_mid = fig.add_subplot(gs[1], sharex=ax_top)
    ax_bot = fig.add_subplot(gs[2], sharex=ax_top)

    ax_top.plot(
        display_df[Columns.TIME],
        display_df[Columns.SBP],
        color="#1B2A4A",
        linewidth=2.4,
        marker="o",
        markersize=5,
        label="SBP",
    )
    ax_top_hr = ax_top.twinx()
    ax_top_hr.plot(
        display_df[Columns.TIME],
        display_df[Columns.HR],
        color="#E67E22",
        linewidth=2.2,
        marker="s",
        markersize=4,
        label="HR",
    )

    ax_mid.plot(
        display_df[Columns.TIME],
        display_df[Columns.DBP],
        color="#2E86AB",
        linewidth=2.6,
        marker="o",
        markersize=5,
        label="Observed DBP",
    )
    ax_mid.plot(
        display_df[Columns.TIME],
        display_df["predicted_dbp"],
        color="#E05263",
        linewidth=2.6,
        linestyle="--",
        label="Predicted DBP from baseline model",
    )

    ax_bot.axhline(0, color="#444444", linewidth=1.4, linestyle="--")
    ax_bot.bar(
        display_df[Columns.TIME],
        display_df["residual_dbp"],
        width=0.0075,
        color="#4E79A7",
        alpha=0.85,
        edgecolor="white",
        linewidth=0.4,
    )

    baseline_span_start = pre_alert_baseline[Columns.TIME].min()
    baseline_span_end = pre_alert_baseline[Columns.TIME].max()
    alert_span_start = alert_df[Columns.TIME].min()
    alert_span_end = alert_df[Columns.TIME].max()
    for axis in (ax_top, ax_mid, ax_bot):
        axis.axvspan(
            baseline_span_start,
            baseline_span_end,
            color="#D6EAF8",
            alpha=0.5,
            zorder=0,
        )
        axis.axvspan(
            alert_span_start,
            alert_span_end,
            color="#FADBD8",
            alpha=0.55,
            zorder=0,
        )
        axis.grid(True, axis="y", alpha=0.22, linewidth=0.8)

    ax_mid.scatter(
        alert_peak[Columns.TIME],
        alert_peak[Columns.DBP],
        s=120,
        facecolors="none",
        edgecolors="#C0392B",
        linewidth=2.2,
        zorder=5,
    )
    ax_top.set_ylabel("SBP (mmHg)", fontsize=13, fontweight="bold", color="#1B2A4A")
    ax_top_hr.set_ylabel("HR (bpm)", fontsize=13, fontweight="bold", color="#E67E22")
    ax_mid.set_ylabel("DBP (mmHg)", fontsize=13, fontweight="bold")
    ax_bot.set_ylabel("DBP residual\n(mmHg)", fontsize=12, fontweight="bold")
    ax_bot.set_xlabel("Clock time", fontsize=13, fontweight="bold")

    for axis in (ax_top, ax_mid, ax_bot, ax_top_hr):
        axis.tick_params(labelsize=11)
    ax_bot.xaxis.set_major_formatter(mdates.DateFormatter("%H:%M"))
    ax_top.tick_params(labelbottom=False)
    ax_mid.tick_params(labelbottom=False)

    ax_top.text(
        baseline_span_start + (baseline_span_end - baseline_span_start) / 2,
        ax_top.get_ylim()[1] * 0.985,
        "Baseline fit window",
        ha="center",
        va="top",
        fontsize=10.5,
        color="#1B4F72",
        fontweight="bold",
    )
    ax_top.text(
        alert_span_start + (alert_span_end - alert_span_start) / 2,
        ax_top.get_ylim()[1] * 0.985,
        "Air-alert window",
        ha="center",
        va="top",
        fontsize=10.5,
        color="#943126",
        fontweight="bold",
    )

    top_handles, top_labels = ax_top.get_legend_handles_labels()
    hr_handles, hr_labels = ax_top_hr.get_legend_handles_labels()
    ax_top.legend(
        top_handles + hr_handles,
        top_labels + hr_labels,
        loc="upper right",
        fontsize=10,
        frameon=True,
        ncol=2,
    )
    ax_mid.legend(loc="upper left", fontsize=10, frameon=True)
    ax_bot.text(
        0.99,
        0.88,
        (
            f"{winner} | baseline n={int(metrics_row['Train_N'])} | alert n={int(metrics_row['DBP_Air Alert_N'])}\n"
            f"baseline MAE {metrics_row['DBP_Ref_MAE']:.2f} mmHg | alert MAE {metrics_row['DBP_Air Alert_MAE']:.2f} mmHg\n"
            f"MAE inflation {metrics_row['DBP_Air Alert_Anomaly']:+.0f}% | bias {metrics_row['DBP_Air Alert_DeltaBias']:+.2f} mmHg"
        ),
        transform=ax_bot.transAxes,
        ha="right",
        va="top",
        fontsize=9.5,
        color="#1F2D3D",
        bbox=dict(boxstyle="round,pad=0.3", fc="white", ec="#CBD5E1", lw=1.0, alpha=0.95),
    )

    fig.subplots_adjust(left=0.08, right=0.92, top=0.94, bottom=0.1)
    fig.savefig(CASE_STUDY_PATH, dpi=250, bbox_inches="tight")
    plt.close(fig)
    return True


def prepare_presentation_assets():
    """Generate presentation-specific figures when the source data are available."""
    try:
        generate_context_breakdown_chart()
    except Exception as exc:
        print(f"Warning: could not regenerate context breakdown figure: {exc}")
    try:
        generate_presentation_cognitive_patterns()
    except Exception as exc:
        print(f"Warning: could not regenerate cognitive presentation figure: {exc}")
    try:
        generate_case_study_timeline()
    except Exception as exc:
        print(f"Warning: could not regenerate case-study figure: {exc}")


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_slide_01_title(prs):
    """Slide 1: Title slide -- dark background, centered content."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY

    # Thin accent stripe near top
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.2), SLIDE_WIDTH, Inches(0.04),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = BLUE_ACCENT
    stripe.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.4),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Detecting Stress-Linked Blood Pressure Relationship Shifts from Ambulatory Monitoring"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Authors
    authors_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(3.1), Inches(11.3), Inches(0.7),
    )
    tf = authors_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "A. Tymchak, A. Butkevych, K. Yurtsiv, I. Nastenko, "
        "D. Zahorska, O. Shaposhnyk, S. Yanushkevich, V. Babenko"
    )
    p.font.size = Pt(16)
    p.font.color.rgb = ICE_BLUE
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Affiliations
    aff_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.9),
    )
    tf = aff_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Igor Sikorsky Kyiv Polytechnic Institute  |  "
        "University of Calgary  |  "
        "Center of Maternity & Childhood, NAMS of Ukraine"
    )
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(0x99, 0xAA, 0xCC)
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Conference
    conf_box = slide.shapes.add_textbox(
        Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.5),
    )
    tf = conf_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IEEE ELNANO 2026  \u2022  Kyiv, Ukraine"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # NATO acknowledgment
    nato_box = slide.shapes.add_textbox(
        Inches(2.5), Inches(5.7), Inches(8.3), Inches(0.5),
    )
    tf = nato_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Supported by NATO Science for Peace and Security Programme \u2013 Grant G8475"
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(0x88, 0x99, 0xBB)
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Bottom stripe
    stripe2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), SLIDE_WIDTH, Inches(0.04),
    )
    stripe2.fill.solid()
    stripe2.fill.fore_color.rgb = BLUE_ACCENT
    stripe2.line.fill.background()

    # No slide number on title slide; add only the IEEE footer
    footer_box = slide.shapes.add_textbox(
        Inches(5.0), Inches(6.95), Inches(3.3), Inches(0.35),
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IEEE ELNANO 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x66, 0x77, 0x99)
    p.font.name = FONT_BODY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide, (
        "Good morning, everyone. My name is Vitalii Babenko, presenting on behalf of our "
        "research team from Igor Sikorsky KPI, the University of Calgary, and the Center of "
        "Maternity and Childhood in Kyiv. Today, we will discuss a computational method "
        "designed to detect stress-linked shifts in the relationship between blood pressure "
        "and heart rate, utilizing standard ambulatory monitoring data. I would like to "
        "gratefully acknowledge the NATO Science for Peace and Security Programme for "
        "supporting this interdisciplinary research. We will examine how shifting from "
        "cohort averages to individualized baseline modeling can help us extract clearer, "
        "contextual stress signals from everyday clinical recordings."
    ))
    return slide


def build_slide_02_motivation(prs):
    """Slide 2: Motivation / Problem."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Motivation: The Specificity Problem")

    bullets = [
        "Stress impacts cardiovascular regulation and blood pressure",
        "Specificity gap in standard Ambulatory Blood Pressure Monitoring (ABPM)",
        "Elevated readings: Acute stress response vs. sustained hypertension",
        "Real-world setting: Kyiv residents facing air-raid alert exposures",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), height=Inches(3.0),
                    font_size=17, bold_first=False)

    # Visual element: key insight callout box (moved up to reduce gap)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(4.0), Inches(10.3), Inches(1.0),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = VERY_LIGHT_BLUE
    box.line.color.rgb = BLUE_ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key question: "
    run.font.bold = True
    run.font.size = Pt(15)
    run.font.color.rgb = NAVY
    run.font.name = FONT_BODY
    run2 = p.add_run()
    run2.text = "Can standard 24-hour ABPM data reveal the physiological footprint of different contexts through changes in hemodynamic coupling?"
    run2.font.size = Pt(15)
    run2.font.color.rgb = CHARCOAL
    run2.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 2)
    add_speaker_notes(slide, (
        "The primary clinical challenge we are addressing is a lack of specificity in standard "
        "cardiovascular monitoring. A 24-hour Ambulatory Blood Pressure Monitor, or ABPM, "
        "successfully captures dynamic changes in blood pressure, but it fundamentally lacks "
        "behavioral context. When a clinician reviews a 24-hour log, an elevated reading could "
        "reflect several different clinical realities. It might be a residual white-coat effect, "
        "emerging sustained hypertension that requires pharmacological intervention, or simply a "
        "normal, acute autonomic response to a transient environmental stressor. Without context, "
        "these highly distinct physiological states look completely identical on paper. This "
        "feasibility study investigates whether we can extract that missing contextual signal "
        "directly from the ABPM data itself. Our dataset was collected in a real-world setting "
        "in Kyiv, where participants navigated their normal daily routines, which unfortunately "
        "included periods of air-raid alert exposure over the winter. We aimed to determine if "
        "these standard, 24-hour recordings could reveal the physiological footprint of different "
        "contexts through measurable changes in an individual's hemodynamic coupling."
    ))
    return slide


def build_slide_03_related_work(prs):
    """Slide 3: Related work gaps."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Related Work Gaps")

    bullets = [
        "Group-level analyses often obscure individual heterogeneity",
        "Traditional monitoring overlooks within-subject hemodynamic coupling",
        "High-fidelity stress detection typically requires specialized equipment",
        "Need for scalable insights from deployed clinical hardware",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), width=Inches(6.0),
                    height=Inches(3.5), font_size=18)

    # Right-side comparison: two stacked cards
    # Card 1: Existing approaches
    card1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(1.5), Inches(4.8), Inches(1.8),
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE
    card1.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card1.line.width = Pt(0.75)

    # Left accent
    acc1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(1.5), Inches(0.08), Inches(1.8),
    )
    acc1.fill.solid()
    acc1.fill.fore_color.rgb = RED_ACCENT
    acc1.line.fill.background()

    tb1 = slide.shapes.add_textbox(Inches(8.1), Inches(1.6), Inches(4.3), Inches(1.6))
    tf = tb1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Existing approaches"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RED_ACCENT
    p.font.name = FONT_BODY
    p2 = tf.add_paragraph()
    p2.text = "Cohort averages, lab-only sensors, HR/HRV confounded by exertion"
    p2.font.size = Pt(13)
    p2.font.color.rgb = CHARCOAL
    p2.font.name = FONT_BODY
    p2.space_before = Pt(6)

    # Card 2: Our approach
    card2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(3.6), Inches(4.8), Inches(1.8),
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card2.line.width = Pt(0.75)

    acc2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.8), Inches(3.6), Inches(0.08), Inches(1.8),
    )
    acc2.fill.solid()
    acc2.fill.fore_color.rgb = BLUE_ACCENT
    acc2.line.fill.background()

    tb2 = slide.shapes.add_textbox(Inches(8.1), Inches(3.7), Inches(4.3), Inches(1.6))
    tf = tb2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Our approach"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = FONT_BODY
    p2 = tf.add_paragraph()
    p2.text = "Within-subject baselines, ABPM-only, software-level coupling deviation screening"
    p2.font.size = Pt(13)
    p2.font.color.rgb = CHARCOAL
    p2.font.name = FONT_BODY
    p2.space_before = Pt(6)

    gap_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9), Inches(4.95), Inches(5.9), Inches(1.0),
    )
    gap_box.fill.solid()
    gap_box.fill.fore_color.rgb = VERY_LIGHT_BLUE
    gap_box.line.color.rgb = BLUE_ACCENT
    gap_box.line.width = Pt(1.2)
    tf = gap_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Missing piece: a scalable way to read individualized physiology from standard ABPM, without lab-only hardware."
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 3)
    add_speaker_notes(slide, (
        "Much of the existing literature regarding cardiovascular stress response relies heavily "
        "on group-level statistical designs, measuring the cohort's average response to a given "
        "stimulus. However, hemodynamic responses to stress are highly heterogeneous. One patient "
        "might react to a stressor primarily with heart rate elevation, while another might "
        "experience profound vasoconstriction with a minimal chronotropic response. Averaging "
        "these reactions across a broad cohort frequently washes out these individualized "
        "physiological adjustments entirely. Furthermore, the studies that do successfully map "
        "out detailed, individual hemodynamic pathways typically utilize specialized, "
        "non-ambulatory equipment such as continuous beat-to-beat finger cuffs or lab-based "
        "impedance cardiography. While valuable, these setups cannot be scaled for prolonged "
        "at-home monitoring. Our objective is to bridge this gap by developing a scalable, "
        "software-level approach that utilizes the widely-available ABPM hardware already "
        "deployed in clinical practice. Instead of comparing patients to a population mean, we "
        "focus on interpreting the internal coupling relationships between a single patient's "
        "own cardiovascular variables over time."
    ))
    return slide


def build_slide_04_data(prs):
    """Slide 4: Data."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Data Collection & Labeling")

    bullets = [
        "28 participants (age 17\u201366; 42.9% female), Kyiv, Nov 2023 \u2013 Feb 2024",
        "2,164 valid readings across 5 labeled context windows",
        "1,238 awake-baseline readings used for training (~44 per participant)",
        "Standard ABPM cadence: one cuff reading every 15\u201330 minutes",
        "8-level deterministic label priority; air-raid alerts take precedence",
        "Cognitive tasks: Stroop, reaction time, visuomotor coordination",
        "Physical task: modified Harvard Step Test",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), width=Inches(5.55),
                    height=Inches(4.0), font_size=17)

    # Right side: context breakdown bar chart — wider, shifted left for more bar room
    try_add_image(
        slide, CONTEXT_BREAKDOWN_PATH,
        x=5.95, y=1.22, w=6.95, h=4.75,
        placeholder_text="[Dataset composition by labeled context]",
    )

    # Bottom note
    note_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(5.6), Inches(6.5), Inches(0.4),
    )
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Training density is modest but sufficient for transparent OLS-style modeling"
    p.font.size = Pt(13)
    p.font.color.rgb = MUTED_GRAY
    p.font.name = FONT_BODY
    p.font.italic = True

    add_footer(slide, 4)
    add_speaker_notes(slide, (
        "The dataset contains 2,164 valid ABPM readings from 28 participants, but the key number "
        "for modeling is 1,238 awake-baseline points. That is about 44 baseline measurements per "
        "person on average. For standard ABPM, this is enough for transparent linear models and "
        "too sparse for data-hungry deep learning. We made that trade-off on purpose. The cuff "
        "inflates every 15 to 30 minutes, so interpretability and variance control matter more "
        "than model complexity here. Contexts were labeled deterministically, with air-raid "
        "alerts at the top of the priority hierarchy, followed by controlled task windows and "
        "baseline at the bottom."
    ))
    return slide


def build_slide_05_methods(prs):
    """Slide 5: Methods -- includes pipeline diagram."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Methods: Individualized Coupling Model")

    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(1.38), Inches(12.0), Inches(1.22),
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = WHITE
    formula_box.line.color.rgb = BLUE_ACCENT
    formula_box.line.width = Pt(1.4)

    formula_text = slide.shapes.add_textbox(
        Inches(0.95), Inches(1.56), Inches(11.4), Inches(0.82),
    )
    tf = formula_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DBP̂ = β0 + Σ βj Φj(x),   |S| ≤ 3"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Φ(x) = [SBP, HR, 1/SBP, 1/HR, SBP·HR, 1/(SBP·HR)]"
    p2.font.size = Pt(20)
    p2.font.color.rgb = CHARCOAL
    p2.font.name = FONT_BODY
    p2.alignment = PP_ALIGN.CENTER

    add_formula_card(
        slide,
        "MAE inflation",
        "Aᵢ,c (%) = 100 × (MAEᵢ(c) - MAEᵢ,ref) / (MAEᵢ,ref + ε)",
        x=0.7, y=2.8, w=5.8, h=0.82,
        accent_color=BLUE_ACCENT,
    )
    add_formula_card(
        slide,
        "DeltaBias",
        "ΔBiasᵢ,c = median(eᵢt),   eᵢt = DBPᵢt - DBP̂ᵢt",
        x=6.7, y=2.8, w=5.8, h=0.82,
        accent_color=RED_ACCENT,
    )

    # Pipeline diagram — expanded to use the freed space
    pipeline_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(3.85), Inches(12.3), Inches(2.8),
    )
    pipeline_bg.fill.solid()
    pipeline_bg.fill.fore_color.rgb = PIPELINE_BG
    pipeline_bg.line.fill.background()

    # Pipeline label
    plabel = slide.shapes.add_textbox(Inches(0.7), Inches(3.98), Inches(3.0), Inches(0.35))
    tf = plabel.text_frame
    p = tf.paragraphs[0]
    p.text = "Analysis Pipeline"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY

    # Pipeline boxes (larger for projection readability)
    box_y = 4.88
    box_h = 1.02
    box_w = 1.72
    gap = 0.14
    arrow_w = 0.28

    steps = [
        "ABPM\nData",
        "Context\nLabeling",
        "Baseline\nModel Fit",
        "Time-Block\nCross-Val",
        "Deviation\nMetrics",
        "Screening\nFlags",
    ]

    start_x = 0.7
    for i, step_text in enumerate(steps):
        bx = start_x + i * (box_w + gap + arrow_w)
        add_pipeline_box(slide, step_text, bx, box_y, box_w, box_h,
                         fill_color=WHITE, text_color=NAVY, font_size=14, bold=True)
        if i < len(steps) - 1:
            ax = bx + box_w + 0.03
            add_arrow(slide, ax, box_y + box_h / 2 - 0.125, arrow_w)

    add_footer(slide, 5)
    add_speaker_notes(slide, (
        "This is the core engineering slide. For each participant, we fit DBP only from that "
        "person's awake baseline readings. The feature pool is explicit: SBP, HR, their inverses, "
        "their product, and the inverse product. We never let the model roam freely through all "
        "six terms. Subset size is capped at three features because the training set is only "
        "about 44 baseline readings per participant on average. That is exactly where curse-of-"
        "dimensionality problems start if you get greedy. Model selection uses three contiguous "
        "time blocks, not shuffled folds. That choice removes temporal leakage, which is the main "
        "failure mode in ABPM time-series modeling. Under this data density, simple linear models "
        "are the right tool: transparent, stable, and sufficient. Once the baseline model is fit, "
        "we quantify stress-window deviation with two outputs only: MAE inflation for magnitude "
        "and signed residual bias for direction."
    ))
    return slide


def build_slide_06_results(prs):
    """Slide 6: Results -- with Figure 1."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Results: Individual-Level Patterns")

    # Prominent null-result banner across the top
    null_banner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(1.2), Inches(11.9), Inches(0.6),
    )
    null_banner.fill.solid()
    null_banner.fill.fore_color.rgb = VERY_LIGHT_BLUE
    null_banner.line.color.rgb = BLUE_ACCENT
    null_banner.line.width = Pt(1.5)
    tf = null_banner.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "Cohort averaging masks the signal (all q > 0.05) \u2192 validating individualized models"
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = NAVY
    run1.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    # Key stats — shifted down to sit below the banner
    add_stat_callout(slide, "6.89", "Baseline MAE\n(mmHg)", 0.7, 2.05, w=2.8, h=1.1,
                     number_color=NAVY)
    add_stat_callout(slide, "+4.18", "Bias for Cognitive\n(mmHg)", 3.8, 2.05, w=2.8, h=1.1,
                     number_color=BLUE_ACCENT)
    add_stat_callout(slide, "-2.01", "Bias for Physical\n(mmHg)", 6.9, 2.05, w=2.8, h=1.1,
                     number_color=RED_ACCENT)
    add_stat_callout(slide, "20/27", "Flagged (74.1%)\n(descriptive threshold)", 10.0, 2.05, w=2.8, h=1.1,
                     number_color=NAVY)

    # Figure 1 -- insert the dot plot
    # Expected: results/fig1.png (two-panel dot plot)
    figure_path = FIG1_PRESENTATION_PATH if os.path.exists(FIG1_PRESENTATION_PATH) else FIG1_PATH
    fig_found = try_add_image(
        slide, figure_path,
        x=1.5, y=3.3, w=10.3, h=3.5,
        placeholder_text=(
            "[Figure 1: Two-panel dot plot showing MAE inflation (Panel A) "
            "and signed residual bias (Panel B) per participant]"
        ),
    )

    # Caption
    cap_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.35),
    )
    tf = cap_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Figure 1: Cognitive-task deviation metrics split by screen status. Threshold lines: 50% inflation and +2 mmHg bias."
    )
    p.font.size = Pt(11)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_BODY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 6)
    add_speaker_notes(slide, (
        "The median baseline error is 6.89 millimeters of mercury. Then we test the coupling "
        "metrics at cohort level, apply BH-FDR, and get a clean null. That is not a failure. It "
        "proves the point: group averaging washes out the physiology we care about. Once we stay "
        "within-subject, the structure is obvious. The left panel shows MAE inflation. The right "
        "panel shows signed residual bias. The red threshold line is the descriptive screening "
        "rule: more than 50 percent error inflation, or more than +2 millimeters of mercury bias. "
        "By that rule, 20 of 27 cognitive-task participants are flagged. The direction also "
        "matters. Cognitive load pushes bias upward. Physical load pushes it downward. So the "
        "residual is not just noise; it carries physiological information that cohort statistics "
        "erase."
    ))
    return slide


def build_slide_07_case_study(prs):
    """Slide 7: Single-subject alert case study."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Case Study: One Patient, One Alert Window")

    try_add_image(
        slide,
        CASE_STUDY_PATH,
        x=0.65, y=1.38, w=9.05, h=5.65,
        placeholder_text="[Single-subject baseline vs air-alert timeline]",
    )

    stat_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.95), Inches(1.52), Inches(2.55), Inches(1.1),
    )
    stat_card.fill.solid()
    stat_card.fill.fore_color.rgb = WHITE
    stat_card.line.color.rgb = BLUE_ACCENT
    stat_card.line.width = Pt(1.2)
    tf = stat_card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Subject 35"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Baseline MAE 5.23 mmHg\nAir-alert MAE 11.82 mmHg"
    p2.font.size = Pt(13)
    p2.font.color.rgb = CHARCOAL
    p2.font.name = FONT_BODY
    p2.alignment = PP_ALIGN.CENTER

    bullets = [
        "38 baseline points | 6 alert points",
        "Winner: OLS(SBP)",
        "Inflation: +126%",
        "Bias: +9.50 mmHg",
    ]
    add_bullet_list(
        slide,
        bullets,
        left=Inches(9.9),
        top=Inches(2.9),
        width=Inches(2.55),
        height=Inches(1.9),
        font_size=15,
    )

    note_box = slide.shapes.add_textbox(
        Inches(9.93), Inches(5.45), Inches(2.45), Inches(1.05),
    )
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "SBP stays within the learned operating range.\nDBP is what breaks away from prediction."
    )
    p.font.size = Pt(12.5)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.name = FONT_BODY
    p.font.italic = True

    add_footer(slide, 7)
    add_speaker_notes(slide, (
        "Here is the patient-level picture that the cohort averages hide. This is Subject 35. "
        "The blue block is the baseline reference window used to learn that person's usual "
        "coupling. The red block is the air-alert window. On the top panel, SBP and heart rate "
        "move, but they stay in a range the baseline model has already seen. The middle panel is "
        "the key result: observed DBP separates from the baseline prediction during the alert. "
        "The highlighted point shows the local breakout. Quantitatively, baseline MAE is 5.23 "
        "millimeters of mercury, while the alert-window MAE jumps to 11.82. That is 126 percent "
        "inflation, with a positive residual bias of 9.50 millimeters of mercury. This is exactly "
        "what we mean by stress-linked coupling disruption at the individual level."
    ))
    return slide


def build_slide_08_discussion(prs):
    """Slide 8: Discussion."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Discussion: Interpreting the Patterns")

    # Two-column layout
    # Left column: bullets
    bullets = [
        "Null group results validate the paradigm: cohort averaging washes out the signal",
        "Residual direction matters: cognitive and physical tasks split in opposite directions",
        "Residual metrics act as a low-cost software proxy for vascular adjustment",
        "Air-alert windows remain heterogeneous, so analysis must stay subject-level",
    ]
    add_bullet_list(slide, bullets, top=Inches(1.4), width=Inches(6.0),
                    height=Inches(3.0), font_size=17)

    # Right column: interpretation cards
    cards = [
        ("Cognitive tasks", "+4.18 mmHg bias", "Consistent with increased\nvascular tone (inferential)", BLUE_ACCENT),
        ("Physical tasks", "-2.01 mmHg bias", "Consistent with exercise-related\nperipheral vasodilation (inferential)", RED_ACCENT),
        ("Air-raid alerts", "+27.35% MAE inflation", "Heterogeneous behavioral\nresponses, near-zero bias", NAVY),
    ]

    for idx, (title, stat, desc, accent_color) in enumerate(cards):
        cy = 1.5 + idx * 1.55  # Even spacing between cards
        # Card bg
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7.5), Inches(cy), Inches(5.1), Inches(1.4),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        card.line.width = Pt(0.75)

        # Left accent
        acc = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7.5), Inches(cy), Inches(0.08), Inches(1.4),
        )
        acc.fill.solid()
        acc.fill.fore_color.rgb = accent_color
        acc.line.fill.background()

        # Title + stat
        tb = slide.shapes.add_textbox(Inches(7.85), Inches(cy + 0.1), Inches(4.5), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = title + "  "
        run1.font.size = Pt(14)
        run1.font.bold = True
        run1.font.color.rgb = accent_color
        run1.font.name = FONT_BODY
        run2 = p.add_run()
        run2.text = stat
        run2.font.size = Pt(14)
        run2.font.color.rgb = CHARCOAL
        run2.font.name = FONT_BODY

        # Description
        db = slide.shapes.add_textbox(Inches(7.85), Inches(cy + 0.55), Inches(4.5), Inches(0.7))
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = FONT_BODY

    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9), Inches(5.0), Inches(5.9), Inches(1.05),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = VERY_LIGHT_BLUE
    callout.line.color.rgb = BLUE_ACCENT
    callout.line.width = Pt(1.2)
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Practical takeaway: keep the analysis subject-level and use residual direction, not cohort means, to interpret the physiology."
    p.font.size = Pt(14.5)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 8)
    add_speaker_notes(slide, (
        "The null group result proves our hypothesis. Cohort averaging washes out vital "
        "physiological signal, so individualized modeling is mandatory here. The sign of the "
        "residual is informative. Cognitive tasks show upward bias, consistent with higher "
        "vascular tone. Physical tasks show downward bias, consistent with exercise-related "
        "vasodilation. Since ABPM does not measure cardiac output or systemic vascular resistance "
        "directly, our residual metrics act as a low-cost software proxy for these vascular "
        "adjustments. That claim remains inferential, but it is physiologically grounded. For "
        "air-raid alerts, the key point is heterogeneity. These are exposure windows, not direct "
        "stress labels. Some participants relocate, some freeze, some keep working. That is why "
        "the alert condition produces variance inflation instead of a single uniform directional "
        "shift."
    ))
    return slide


def build_slide_09_limitations(prs):
    """Slide 9: Boundary conditions and future work."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Boundary Conditions of the Feasibility Study")

    # Two-column layout
    # Left: Limitations
    lim_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.3), Inches(5.5), Inches(0.4),
    )
    tf = lim_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Feasibility Boundary Conditions"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY

    lim_bullets = [
        "N = 28 is strong for feasibility, but still narrow for subgroup claims",
        "ABPM cadence stays sparse at 15\u201330 min, so transients are under-sampled",
        "Lack of continuous cardiac output ground truth",
        "Coarse context labels do not capture moment-to-moment arousal",
    ]
    add_bullet_list(slide, lim_bullets, left=Inches(0.7), top=Inches(1.85),
                    width=Inches(5.5), height=Inches(2.5), font_size=15)

    # Right: Future Work
    fut_title = slide.shapes.add_textbox(
        Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.4),
    )
    tf = fut_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Validation Steps"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_ACCENT
    p.font.name = FONT_BODY

    fut_bullets = [
        "Larger cohorts with higher-resolution ambulatory monitors",
        "Impedance cardiography calibration sub-study",
        "Integration with cuffless / wearable BP platforms",
        "Ecological momentary assessment prompts at cuff inflation",
    ]
    add_bullet_list(slide, fut_bullets, left=Inches(7.0), top=Inches(1.85),
                    width=Inches(5.8), height=Inches(3.0), font_size=15)

    # Visual separator line (more visible)
    sep = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.55), Inches(1.5), Inches(0.04), Inches(3.5),
    )
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
    sep.line.fill.background()

    # Bottom callout (more generous margins from slide edges)
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.8), Inches(5.2), Inches(9.7), Inches(1.0),
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = VERY_LIGHT_BLUE
    callout.line.color.rgb = BLUE_ACCENT
    callout.line.width = Pt(1.5)
    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "This dataset is strong enough to establish the method. The next step is calibration "
        "and external validation, not retreat from the signal."
    )
    p.font.size = Pt(14)
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, 9)
    add_speaker_notes(slide, (
        "These are boundary conditions, not apologies. Clean 24-hour ABPM data with controlled "
        "tasks, naturalistic alerts, and synchronized labels are hard to collect, especially in "
        "wartime Kyiv. The sample is fully adequate for a feasibility study and for establishing "
        "the pipeline. But it is still too small for confident subgroup claims, especially in the "
        "air-alert branch. The second constraint is sampling density: a cuff reading every 15 to "
        "30 minutes will miss fast autonomic transitions. The third is mechanistic calibration: "
        "without continuous cardiac output or SVR, physiology remains inferential. So the next "
        "move is straightforward. Replicate on larger cohorts, add higher-resolution wearable "
        "sensors, and run an impedance-cardiography sub-study to anchor the residual metrics to "
        "hemodynamic ground truth."
    ))
    return slide


def build_slide_10_conclusions(prs):
    """Slide 10: Conclusions."""
    slide = add_blank_slide(prs)
    add_header_bar(slide, "Conclusions")

    conclusions = [
        ("Standard ABPM data yields within-subject stress screening signals",
         "24-hour recordings contain context-sensitive information when analyzed at the individual level"),
        ("Individualized baselines capture context missed by cohort averages",
         "Moving from group-average to per-subject models reveals regulatory decoupling episodes"),
        ("Dual-metric approach distinguishes magnitude and direction of coupling deviation",
         "MAE inflation quantifies shift magnitude; signed bias classifies hemodynamic patterns"),
        ("Scalable, software-level analysis for behavioral cardiovascular monitoring",
         "Computationally lightweight method to increase specificity in longitudinal screening"),
    ]

    y_start = 1.5
    for idx, (headline, detail) in enumerate(conclusions):
        cy = y_start + idx * 1.3

        # Number circle
        circ = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(0.9), Inches(cy + 0.05), Inches(0.45), Inches(0.45),
        )
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE_ACCENT
        circ.line.fill.background()
        tf = circ.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(idx + 1)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_HEADER
        p.alignment = PP_ALIGN.CENTER

        # Headline
        hb = slide.shapes.add_textbox(
            Inches(1.6), Inches(cy), Inches(10.5), Inches(0.4),
        )
        tf = hb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = headline
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.font.name = FONT_BODY

        # Detail
        db = slide.shapes.add_textbox(
            Inches(1.6), Inches(cy + 0.42), Inches(10.5), Inches(0.5),
        )
        tf = db.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(13)
        p.font.color.rgb = MEDIUM_GRAY
        p.font.name = FONT_BODY

    add_footer(slide, 10)
    add_speaker_notes(slide, (
        "Standard ABPM already contains subject-level stress physiology. The problem is not the "
        "sensor. The problem is the analysis level. Once we stop averaging people together and "
        "start modeling each participant against their own baseline, coupling disruptions become "
        "visible. Two numbers are enough to screen them: MAE inflation for magnitude and signed "
        "bias for direction. That gives us a lightweight, interpretable path toward context-aware "
        "cardiovascular monitoring."
    ))
    return slide


def build_slide_11_thank_you(prs):
    """Slide 11: Acknowledgments and Thank You."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_NAVY

    # Top accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.0), SLIDE_WIDTH, Inches(0.04),
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = BLUE_ACCENT
    stripe.line.fill.background()

    # "Thank You" title
    ty_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(1.4), Inches(11.3), Inches(1.0),
    )
    tf = ty_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Acknowledgments
    ack_items = [
        "Funding: NATO Science for Peace and Security Programme (Grant G8475)",
        "Collaborators: Igor Sikorsky KPI & University of Calgary",
        "Special thanks to the study participants in Kyiv",
    ]
    ack_box = slide.shapes.add_textbox(
        Inches(2.0), Inches(2.8), Inches(9.3), Inches(2.2),
    )
    tf = ack_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(ack_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = ICE_BLUE
        p.font.name = FONT_BODY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(12)

    # Questions box
    q_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.5), Inches(5.2), Inches(4.3), Inches(0.8),
    )
    q_box.fill.solid()
    q_box.fill.fore_color.rgb = BLUE_ACCENT
    q_box.line.fill.background()
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_HEADER
    p.alignment = PP_ALIGN.CENTER

    # Bottom stripe
    stripe2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.6), SLIDE_WIDTH, Inches(0.04),
    )
    stripe2.fill.solid()
    stripe2.fill.fore_color.rgb = BLUE_ACCENT
    stripe2.line.fill.background()

    # Footer
    footer_box = slide.shapes.add_textbox(
        Inches(5.0), Inches(6.85), Inches(3.3), Inches(0.35),
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "IEEE ELNANO 2026"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x66, 0x77, 0x99)
    p.font.name = FONT_BODY
    p.font.italic = True
    p.alignment = PP_ALIGN.CENTER

    # Slide number
    num_box = slide.shapes.add_textbox(
        Inches(12.1), Inches(6.85), Inches(0.9), Inches(0.35),
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = "11"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x66, 0x77, 0x99)
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.RIGHT

    add_speaker_notes(slide, (
        "I would like to conclude by formally thanking the NATO Science for Peace and Security "
        "Programme under grant G8475 for funding this research. I also extend my gratitude to "
        "my collaborators, and most importantly, to our study participants for their remarkable "
        "adherence to the clinical protocol while navigating challenging environmental conditions "
        "in Kyiv. Thank you very much for your time and attention today. I would now be happy to "
        "take any questions you may have."
    ))
    return slide


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    prepare_presentation_assets()
    prs = create_presentation()

    build_slide_01_title(prs)
    build_slide_02_motivation(prs)
    build_slide_03_related_work(prs)
    build_slide_04_data(prs)
    build_slide_05_methods(prs)
    build_slide_06_results(prs)
    build_slide_07_case_study(prs)
    build_slide_08_discussion(prs)
    build_slide_09_limitations(prs)
    build_slide_10_conclusions(prs)
    build_slide_11_thank_you(prs)

    output_path = os.path.join(SCRIPT_DIR, "presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
